#!/usr/bin/env python3
"""
Generate the eddacraft/Anvil 13-slide investor pitch deck.

Uses the Nordic Terminal design system:
- Background: --void (#0d0d0f)
- Surface: --surface (#141416)
- Structure: --structure (#2a2a2e)
- Text primary: --text-primary (#ebebeb)
- Text muted: --text-muted (#85858a)
- Anvil accent: --anvil (#cc5500)
- Growth accent: --edda (#2e8b57)
- Error: #c94a4a
- Warning: #d08c38

Typography: JetBrains Mono for headlines/data, Inter for body.
All sharp corners, no shadows, no gradients.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────────────────────

VOID = RGBColor(0x0D, 0x0D, 0x0F)
SURFACE = RGBColor(0x14, 0x14, 0x16)
STRUCTURE = RGBColor(0x2A, 0x2A, 0x2E)
TEXT_PRIMARY = RGBColor(0xEB, 0xEB, 0xEB)
TEXT_MUTED = RGBColor(0x85, 0x85, 0x8A)
ANVIL = RGBColor(0xCC, 0x55, 0x00)
EDDA = RGBColor(0x2E, 0x8B, 0x57)
ERROR = RGBColor(0xC9, 0x4A, 0x4A)
WARNING = RGBColor(0xD0, 0x8C, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ANVIL_DIM = RGBColor(0x33, 0x19, 0x05)  # ~15% opacity anvil on void

# ── Dimensions ───────────────────────────────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_W = SLIDE_WIDTH - 2 * MARGIN

# ── Helpers ──────────────────────────────────────────────────────────────────

def set_slide_bg(slide, color=VOID):
    """Set solid background colour on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name="Inter",
                font_size=16, color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP):
    """Add a text box with single run styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_name="Inter",
                          font_size=16, color=TEXT_PRIMARY, bold=False,
                          alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Clear default paragraph
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        if isinstance(line, tuple):
            # (text, font_name, font_size, color, bold)
            txt, fn, fs, c, b = line
            run = p.add_run()
            run.text = txt
            run.font.name = fn
            run.font.size = Pt(fs)
            run.font.color.rgb = c
            run.font.bold = b
        elif isinstance(line, list):
            # List of (text, font_name, font_size, color, bold) runs in one paragraph
            for run_spec in line:
                txt, fn, fs, c, b = run_spec
                run = p.add_run()
                run.text = txt
                run.font.name = fn
                run.font.size = Pt(fs)
                run.font.color.rgb = c
                run.font.bold = b
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.bold = bold
    return txBox


def add_rect(slide, left, top, width, height, fill_color=None, border_color=STRUCTURE,
             border_width=Pt(1)):
    """Add a rectangle shape with optional fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.color.rgb = border_color
    shape.line.width = border_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    # Rectangles have sharp corners by default (no radius adjustment needed)
    return shape


def add_line(slide, start_x, start_y, end_x, end_y, color=STRUCTURE, width=Pt(1)):
    """Add a line connector."""
    connector = slide.shapes.add_connector(
        1,  # straight connector
        start_x, start_y, end_x, end_y
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_footer(slide, text="[ \u25a0 ] e d d a c r a f t"):
    """Add the eddacraft footer to every slide."""
    add_textbox(
        slide, MARGIN, SLIDE_HEIGHT - Inches(0.5),
        CONTENT_W, Inches(0.4),
        text,
        font_name="JetBrains Mono", font_size=10, color=TEXT_MUTED,
        alignment=PP_ALIGN.LEFT
    )


def add_presenter_notes(slide, notes_text):
    """Add presenter notes to a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_horizontal_bar(slide, left, top, width, height, fill_color, label_text="",
                       label_color=TEXT_MUTED, value_text="", value_color=ANVIL):
    """Add a horizontal bar with optional label and value."""
    bar = add_rect(slide, left, top, width, height, fill_color=fill_color,
                   border_color=fill_color, border_width=Pt(0))
    if label_text:
        add_textbox(slide, left - Inches(2.2), top - Pt(2), Inches(2.1), height,
                    label_text, font_name="JetBrains Mono", font_size=12,
                    color=label_color, alignment=PP_ALIGN.RIGHT)
    if value_text:
        add_textbox(slide, left + width + Inches(0.1), top - Pt(2), Inches(1), height,
                    value_text, font_name="JetBrains Mono", font_size=14,
                    color=value_color, bold=True)
    return bar


# ── Slide Builders ───────────────────────────────────────────────────────────

def build_slide_01_title(prs):
    """Slide 1: Title — AI governance for developers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide)

    # Anvil macro logo as text art in EMBER
    logo_text = (
        "\u2588\u2588\u2588\u2588     \u2588\u2588\u2588\u2588\n"
        "\u2588\u2588         \u2588\u2588\n"
        "\u2588\u2588  \u2588\u2588\u2588\u2588\u2588  \u2588\u2588\n"
        "\u2588\u2588         \u2588\u2588\n"
        "\u2588\u2588  \u2588\u2588\u2588\u2588\u2588  \u2588\u2588\n"
        "\u2588\u2588         \u2588\u2588\n"
        "\u2588\u2588\u2588\u2588     \u2588\u2588\u2588\u2588"
    )
    add_textbox(slide, Inches(4.5), Inches(1.0), Inches(4.3), Inches(2.8),
                logo_text, font_name="JetBrains Mono", font_size=20, color=ANVIL,
                alignment=PP_ALIGN.CENTER)

    # "a n v i l" text to the right of logo
    add_textbox(slide, Inches(8.5), Inches(2.2), Inches(3), Inches(0.5),
                "a n v i l", font_name="JetBrains Mono", font_size=24, color=TEXT_PRIMARY,
                alignment=PP_ALIGN.LEFT)

    # Headline
    add_textbox(slide, MARGIN, Inches(4.2), CONTENT_W, Inches(0.8),
                "AI governance for developers",
                font_name="JetBrains Mono", font_size=44, color=TEXT_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    # Tagline
    add_textbox(slide, MARGIN, Inches(5.0), CONTENT_W, Inches(0.6),
                "Anvil by eddacraft",
                font_name="Inter", font_size=24, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Capabilities
    caps = [
        "Deterministic policy enforcement at file save",
        "Line-level authorship attribution",
        "Architecture drift detection",
        "Policy-as-code via OPA/Rego",
    ]
    add_multiline_textbox(slide, Inches(3.5), Inches(5.7), Inches(6.3), Inches(1.3),
                          caps, font_name="Inter", font_size=14, color=TEXT_MUTED,
                          alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "Anvil is a governance engine for AI-assisted codebases. It enforces policy at the moment code is generated -- at file save -- not after commit. Think of it as the constitutional layer for your repository. We will show you why this matters now, how it works, and why no other tool does what Anvil does.")


def build_slide_02_problem(prs):
    """Slide 2: The Problem — AI writes half the code."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    # Left side — text (40%)
    add_textbox(slide, MARGIN, Inches(0.6), Inches(4.8), Inches(0.8),
                "AI writes half the code.\nNobody governs it.",
                font_name="JetBrains Mono", font_size=36, color=TEXT_PRIMARY, bold=True)

    bullets = [
        "46% of production code is now AI-generated",
        "Fewer than half of developers review AI output",
        "45% of AI-generated code fails security tests",
        "Code duplication increased 4x since AI adoption",
    ]
    add_multiline_textbox(slide, MARGIN, Inches(2.0), Inches(4.5), Inches(2.5),
                          bullets, font_name="Inter", font_size=15, color=TEXT_MUTED)

    # Data callout
    add_textbox(slide, MARGIN, Inches(4.8), Inches(4.5), Inches(1.0),
                "1.7x",
                font_name="JetBrains Mono", font_size=72, color=ANVIL, bold=True)
    add_textbox(slide, MARGIN, Inches(5.8), Inches(4.5), Inches(0.4),
                "defect multiplier for AI-generated code",
                font_name="Inter", font_size=14, color=TEXT_MUTED)

    # Right side — horizontal bar chart (60%)
    chart_left = Inches(5.8)
    chart_w_max = Inches(6.5)
    bar_h = Inches(0.35)
    row_h = Inches(1.05)
    chart_top = Inches(0.9)

    # Chart container
    add_rect(slide, chart_left - Inches(0.1), chart_top - Inches(0.2),
             chart_w_max + Inches(0.5), Inches(5.6), fill_color=SURFACE)

    metrics = [
        ("Issues / PR", 6.45, 10.83, "10.83", "6.45", "1.7x"),
        ("Critical", 1.0, 1.4, "1.4x", "1.0x", "1.4x"),
        ("Major", 1.0, 1.7, "1.7x", "1.0x", "1.7x"),
        ("Security", 1.0, 1.57, "1.57x", "1.0x", "1.57x"),
        ("Maintainability", 1.0, 1.64, "1.64x", "1.0x", "1.64x"),
    ]

    # Normalise to max value for bar widths
    max_val = 10.83
    bar_area_w = Inches(4.8)

    for i, (label, human_val, ai_val, ai_label, human_label, mult) in enumerate(metrics):
        y = chart_top + Inches(0.15) + i * row_h

        # Label
        add_textbox(slide, chart_left, y, Inches(1.8), Inches(0.3),
                    label, font_name="JetBrains Mono", font_size=12, color=TEXT_MUTED)

        bar_left = chart_left + Inches(1.9)
        # Human bar
        human_w = int(bar_area_w * (human_val / max_val))
        add_rect(slide, bar_left, y + Inches(0.3), human_w, bar_h,
                 fill_color=TEXT_MUTED, border_color=TEXT_MUTED, border_width=Pt(0))
        add_textbox(slide, bar_left + human_w + Inches(0.05), y + Inches(0.25),
                    Inches(0.8), bar_h, human_label,
                    font_name="JetBrains Mono", font_size=10, color=TEXT_MUTED)

        # AI bar
        ai_w = int(bar_area_w * (ai_val / max_val))
        add_rect(slide, bar_left, y + Inches(0.68), ai_w, bar_h,
                 fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))
        add_textbox(slide, bar_left + ai_w + Inches(0.05), y + Inches(0.63),
                    Inches(0.8), bar_h, ai_label,
                    font_name="JetBrains Mono", font_size=10, color=ANVIL)

    # Legend
    add_rect(slide, chart_left + Inches(0.1), chart_top + Inches(5.15),
             Inches(0.3), Inches(0.2), fill_color=TEXT_MUTED, border_color=TEXT_MUTED, border_width=Pt(0))
    add_textbox(slide, chart_left + Inches(0.5), chart_top + Inches(5.1),
                Inches(1.0), Inches(0.3), "Human",
                font_name="JetBrains Mono", font_size=10, color=TEXT_MUTED)

    add_rect(slide, chart_left + Inches(1.8), chart_top + Inches(5.15),
             Inches(0.3), Inches(0.2), fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))
    add_textbox(slide, chart_left + Inches(2.2), chart_top + Inches(5.1),
                Inches(1.0), Inches(0.3), "AI-generated",
                font_name="JetBrains Mono", font_size=10, color=ANVIL)

    # Source
    add_textbox(slide, chart_left, chart_top + Inches(5.35), chart_w_max, Inches(0.3),
                "Source: CodeRabbit, State of AI vs Human Code Generation (2025)",
                font_name="Inter", font_size=9, color=TEXT_MUTED)

    add_footer(slide)
    add_presenter_notes(slide, "AI coding tools are mainstream -- 84% of developers use them, 90% of the Fortune 100 have adopted Copilot. But the data is clear: AI-generated code is measurably lower quality. CodeRabbit analysed thousands of pull requests and found 1.7 times more issues in AI-generated PRs. GitClear found code duplication has quadrupled. The productivity gains are real, but so is the quality crisis. And here is the structural problem: every governance tool in the market scans after commit. By then, the ungoverned code is already in the codebase.")


def build_slide_03_why_now(prs):
    """Slide 3: Why Now — The compliance clock is ticking."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.5), CONTENT_W, Inches(0.8),
                "The compliance clock is ticking",
                font_name="JetBrains Mono", font_size=40, color=TEXT_PRIMARY, bold=True)

    # Data callout
    add_textbox(slide, MARGIN, Inches(1.3), Inches(5), Inches(0.8),
                "August 2026",
                font_name="JetBrains Mono", font_size=64, color=ANVIL, bold=True)
    add_textbox(slide, MARGIN, Inches(2.2), Inches(5), Inches(0.4),
                "EU AI Act high-risk enforcement — 5 months away",
                font_name="Inter", font_size=16, color=TEXT_MUTED)

    # Key stats
    stats = [
        "Penalties: up to 7% of global annual turnover",
        "Gartner: 40% of AI coding projects cancelled by 2027",
        "75% of tech leaders face moderate to severe AI technical debt",
    ]
    add_multiline_textbox(slide, MARGIN, Inches(2.9), Inches(5), Inches(1.5),
                          stats, font_name="Inter", font_size=15, color=TEXT_PRIMARY)

    # Timeline
    timeline_top = Inches(4.8)
    timeline_left = MARGIN
    timeline_right = SLIDE_WIDTH - MARGIN
    timeline_w = timeline_right - timeline_left

    # Timeline track
    add_line(slide, timeline_left, timeline_top + Inches(0.3),
             timeline_right, timeline_top + Inches(0.3), color=STRUCTURE, width=Pt(2))

    milestones = [
        ("Feb 2025", "Prohibitions\nactive", TEXT_MUTED, 0.0, 10),
        ("Aug 2025", "GPAI\nobligations", TEXT_MUTED, 0.2, 10),
        ("Mar 2026", "NOW", ANVIL, 0.4, 14),
        ("Aug 2026", "High-risk\nenforcement", ANVIL, 0.58, 14),
        ("2027", "Full\nenforcement", TEXT_MUTED, 0.75, 10),
        ("2030", ">$1B\ngovernance", EDDA, 0.95, 10),
    ]

    for label, desc, color, pos, dot_size in milestones:
        x = int(timeline_left + timeline_w * pos)
        # Dot
        dot_sz = Inches(dot_size / 72)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x - dot_sz // 2,
                                     timeline_top + Inches(0.3) - dot_sz // 2,
                                     dot_sz, dot_sz)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        # Date label above
        add_textbox(slide, x - Inches(0.5), timeline_top - Inches(0.15),
                    Inches(1.0), Inches(0.3), label,
                    font_name="JetBrains Mono", font_size=10, color=color,
                    alignment=PP_ALIGN.CENTER, bold=(color == ANVIL))
        # Description below
        add_textbox(slide, x - Inches(0.6), timeline_top + Inches(0.55),
                    Inches(1.2), Inches(0.6), desc,
                    font_name="Inter", font_size=10, color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)

    # Spend forecast bars (upper right)
    bar_left = Inches(8.5)
    bar_top = Inches(1.0)
    add_rect(slide, bar_left - Inches(0.2), bar_top - Inches(0.2),
             Inches(4.3), Inches(3.2), fill_color=SURFACE)

    add_textbox(slide, bar_left, bar_top, Inches(3.8), Inches(0.4),
                "AI Governance Platform Spend",
                font_name="JetBrains Mono", font_size=14, color=TEXT_PRIMARY)

    # 2026 bar
    add_rect(slide, bar_left, bar_top + Inches(0.7), Inches(2.0), Inches(0.6),
             fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))
    add_textbox(slide, bar_left + Inches(0.1), bar_top + Inches(0.75),
                Inches(1.8), Inches(0.5), "USD 492M",
                font_name="JetBrains Mono", font_size=16, color=VOID, bold=True)
    add_textbox(slide, bar_left + Inches(2.1), bar_top + Inches(0.75),
                Inches(1.0), Inches(0.5), "2026",
                font_name="JetBrains Mono", font_size=12, color=TEXT_MUTED)

    # 2030 bar
    add_rect(slide, bar_left, bar_top + Inches(1.6), Inches(3.8), Inches(0.6),
             fill_color=EDDA, border_color=EDDA, border_width=Pt(0))
    add_textbox(slide, bar_left + Inches(0.1), bar_top + Inches(1.65),
                Inches(2.0), Inches(0.5), ">USD 1B",
                font_name="JetBrains Mono", font_size=16, color=VOID, bold=True)
    add_textbox(slide, bar_left + Inches(3.9), bar_top + Inches(1.65),
                Inches(1.0), Inches(0.5), "2030",
                font_name="JetBrains Mono", font_size=12, color=TEXT_MUTED)

    add_textbox(slide, bar_left, bar_top + Inches(2.6), Inches(3.8), Inches(0.3),
                "Gartner, Feb 2026",
                font_name="Inter", font_size=9, color=TEXT_MUTED)

    add_footer(slide)
    add_presenter_notes(slide, "Three forces are converging. First, regulatory deadlines -- the EU AI Act high-risk requirements become enforceable in August 2026. That is five months from now. Non-compliance penalties reach 7% of global turnover. Second, market reality -- Gartner predicts 40% of AI coding projects will be cancelled by 2027 due to escalating costs and weak governance. Third, budget creation -- Gartner forecasts AI governance platform spend at nearly half a billion dollars this year, growing past a billion by 2030. The market is moving from 'should we govern AI code?' to 'how do we govern AI code?' We have the answer.")


def build_slide_04_solution(prs):
    """Slide 4: Solution — Deterministic governance at file save."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.8), CONTENT_W, Inches(0.8),
                "Deterministic governance at file save",
                font_name="JetBrains Mono", font_size=40, color=TEXT_PRIMARY,
                bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(0.5),
                "Pre-commit — the only governance tool at this position in the workflow",
                font_name="Inter", font_size=16, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Capability card
    card_left = Inches(2.5)
    card_w = Inches(8.3)
    card_top = Inches(2.5)
    card_h = Inches(3.8)
    add_rect(slide, card_left, card_top, card_w, card_h,
             fill_color=SURFACE, border_color=STRUCTURE)

    capabilities = [
        ("[ = ]", "  Policy enforcement at file save — not after commit"),
        ("[ = ]", "  Deterministic analysis — not AI reviewing AI"),
        ("[ = ]", "  Line-level authorship: human / AI / mixed / unknown"),
        ("[ = ]", "  Architecture drift detection via semantic graph"),
        ("[ = ]", "  Policy-as-code (OPA/Rego) — your team controls the rules"),
    ]

    for i, (prefix, desc) in enumerate(capabilities):
        y = card_top + Inches(0.4) + i * Inches(0.62)
        # Prefix in anvil
        add_textbox(slide, card_left + Inches(0.5), y, Inches(0.8), Inches(0.4),
                    prefix, font_name="JetBrains Mono", font_size=18, color=ANVIL,
                    bold=True)
        # Description
        add_textbox(slide, card_left + Inches(1.3), y, Inches(6.5), Inches(0.4),
                    desc, font_name="Inter", font_size=18, color=TEXT_PRIMARY)

    add_footer(slide)
    add_presenter_notes(slide, "Anvil enforces governance at file save -- the moment code is generated. This is architecturally different from every other tool in the market. Static analysers scan after commit. AI review tools evaluate at PR time. Anvil operates at generation time. And critically, Anvil is deterministic. It uses policy-as-code -- OPA and Rego -- not another AI model. The same input always produces the same output. No probabilistic uncertainty. No AI reviewing AI. Every line of code is classified: human, AI, mixed, or unknown. And the architecture of the codebase is tracked incrementally, so drift is detected as a trajectory, not just a violation.")


def build_slide_05_how_it_works(prs):
    """Slide 5: How It Works — Pipeline flow diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.7),
                "File save to governance event in milliseconds",
                font_name="JetBrains Mono", font_size=36, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Pipeline nodes
    stages = [
        ("File Save", "Developer saves file"),
        ("Parse", "Tree-sitter\nincremental (Rust)"),
        ("Attribute", "Line-level\nauthorship"),
        ("Evaluate", "OPA/Rego\npolicy check"),
        ("Govern", "Emit governance\nevent"),
    ]

    node_w = Inches(2.0)
    node_h = Inches(1.4)
    gap = Inches(0.35)
    total_w = 5 * node_w + 4 * gap
    start_x = (SLIDE_WIDTH - total_w) // 2
    node_y = Inches(1.6)

    for i, (label, desc) in enumerate(stages):
        x = start_x + i * (node_w + gap)

        # Node background
        add_rect(slide, x, node_y, node_w, node_h,
                 fill_color=SURFACE, border_color=STRUCTURE)

        # Left accent bar
        add_rect(slide, x, node_y, Pt(4), node_h,
                 fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))

        # Stage label
        add_textbox(slide, x + Inches(0.15), node_y + Inches(0.15),
                    node_w - Inches(0.3), Inches(0.4), label,
                    font_name="JetBrains Mono", font_size=16, color=TEXT_PRIMARY,
                    bold=True, alignment=PP_ALIGN.CENTER)

        # Description
        add_textbox(slide, x + Inches(0.1), node_y + Inches(0.55),
                    node_w - Inches(0.2), Inches(0.7), desc,
                    font_name="Inter", font_size=12, color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)

        # Arrow to next
        if i < 4:
            arrow_x = x + node_w
            arrow_y = node_y + node_h // 2
            # Simple line arrow
            add_line(slide, arrow_x + Inches(0.02), arrow_y,
                     arrow_x + gap - Inches(0.02), arrow_y,
                     color=TEXT_MUTED, width=Pt(2))
            # Arrowhead triangle
            tri = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE,
                arrow_x + gap - Inches(0.15), arrow_y - Inches(0.06),
                Inches(0.12), Inches(0.12)
            )
            tri.fill.solid()
            tri.fill.fore_color.rgb = ANVIL
            tri.line.fill.background()
            tri.rotation = 90.0

    # Output states
    output_y = Inches(3.5)
    add_textbox(slide, MARGIN, output_y, CONTENT_W, Inches(0.4),
                "Governance Output States",
                font_name="JetBrains Mono", font_size=14, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    outputs = [
        ("PASS", "Policy met. Architecture stable.", EDDA),
        ("WARN", "Boundary stress increasing. Review.", WARNING),
        ("BLOCK", "Trust invariant violated. Fix required.", ERROR),
    ]

    out_left = Inches(2.5)
    for i, (state, desc, color) in enumerate(outputs):
        y = output_y + Inches(0.5) + i * Inches(0.55)

        # State card
        add_rect(slide, out_left, y, Inches(8.3), Inches(0.45),
                 fill_color=SURFACE, border_color=STRUCTURE)

        # Left accent
        add_rect(slide, out_left, y, Pt(4), Inches(0.45),
                 fill_color=color, border_color=color, border_width=Pt(0))

        add_textbox(slide, out_left + Inches(0.2), y + Inches(0.05),
                    Inches(1.0), Inches(0.35), state,
                    font_name="JetBrains Mono", font_size=14, color=color, bold=True)

        add_textbox(slide, out_left + Inches(1.3), y + Inches(0.05),
                    Inches(6.5), Inches(0.35), desc,
                    font_name="Inter", font_size=14, color=TEXT_PRIMARY)

    # Terminal command preview
    add_textbox(slide, MARGIN, Inches(5.8), CONTENT_W, Inches(0.5),
                "save \u2192 parse \u2192 attribute \u2192 evaluate \u2192 govern",
                font_name="JetBrains Mono", font_size=18, color=ANVIL,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(6.3), CONTENT_W, Inches(0.4),
                "The entire loop runs in milliseconds. No workflow disruption.",
                font_name="Inter", font_size=14, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "Here is the technical flow. When a developer saves a file, Anvil parses it incrementally using tree-sitter in Rust -- fast enough to be synchronous. Each line is attributed: was this written by a human, an AI assistant, or some combination? Then the policy engine evaluates the change against your team's rules -- these are standard OPA/Rego policies, not proprietary. The architecture graph updates: has this change introduced a new dependency? Crossed a boundary? Expanded the trust surface? Finally, Anvil emits a governance event -- pass, warn, or block. The entire loop runs in milliseconds. No workflow disruption.")


def build_slide_06_product(prs):
    """Slide 6: Product — Built in Rust. No AI inside. 50ms per check."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                "Built in Rust. No AI inside. 50ms per check.",
                font_name="JetBrains Mono", font_size=32, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # TUI mockup frame
    frame_left = Inches(0.5)
    frame_top = Inches(1.1)
    frame_w = Inches(12.3)
    frame_h = Inches(5.5)

    # Outer frame
    add_rect(slide, frame_left, frame_top, frame_w, frame_h,
             fill_color=VOID, border_color=STRUCTURE, border_width=Pt(2))

    # Header bar
    header_h = Inches(0.8)
    add_rect(slide, frame_left, frame_top, frame_w, header_h,
             fill_color=SURFACE, border_color=STRUCTURE)

    # Anvil logo in header
    add_textbox(slide, frame_left + Inches(0.3), frame_top + Inches(0.1),
                Inches(4), Inches(0.6),
                "\u2588\u2588\u2588\u2588  \u2588\u2588\u2588\u2588  a n v i l",
                font_name="JetBrains Mono", font_size=14, color=ANVIL)

    add_textbox(slide, frame_left + Inches(8), frame_top + Inches(0.2),
                Inches(4), Inches(0.4),
                "GOVERNANCE WATCHER  |  POLICY: active  |  50ms",
                font_name="JetBrains Mono", font_size=11, color=TEXT_MUTED)

    # Left pane — policy
    pane_top = frame_top + header_h
    left_pane_w = Inches(4.5)
    add_rect(slide, frame_left, pane_top, left_pane_w, frame_h - header_h - Inches(0.6),
             fill_color=SURFACE, border_color=STRUCTURE)

    add_textbox(slide, frame_left + Inches(0.2), pane_top + Inches(0.1),
                left_pane_w - Inches(0.4), Inches(0.3),
                "[ \u2261 ] ACTIVE_POLICY",
                font_name="JetBrains Mono", font_size=12, color=ANVIL)

    policy_lines = [
        "package anvil.governance",
        "",
        "default allow = false",
        "",
        'allow if {',
        '    input.authorship != "unknown"',
        '    input.trust_score >= 0.7',
        '    not boundary_violation',
        '}',
        "",
        "boundary_violation if {",
        '    input.crosses_boundary',
        '    not input.has_exception',
        "}",
    ]
    add_multiline_textbox(slide, frame_left + Inches(0.3), pane_top + Inches(0.5),
                          left_pane_w - Inches(0.6), Inches(3.5),
                          policy_lines, font_name="JetBrains Mono", font_size=11,
                          color=TEXT_PRIMARY, line_spacing=1.2)

    # Right pane — events
    right_pane_left = frame_left + left_pane_w
    right_pane_w = frame_w - left_pane_w
    add_rect(slide, right_pane_left, pane_top, right_pane_w, frame_h - header_h - Inches(0.6),
             fill_color=VOID, border_color=STRUCTURE)

    add_textbox(slide, right_pane_left + Inches(0.2), pane_top + Inches(0.1),
                right_pane_w - Inches(0.4), Inches(0.3),
                "[ = ] SIGNAL_INTERCEPTOR",
                font_name="JetBrains Mono", font_size=12, color=ANVIL)

    events = [
        [("\u2588 PASS ", "JetBrains Mono", 11, EDDA, True),
         (" src/auth/login.ts  authorship:human  trust:0.92", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 PASS ", "JetBrains Mono", 11, EDDA, True),
         (" src/api/routes.ts  authorship:human  trust:0.88", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 WARN ", "JetBrains Mono", 11, WARNING, True),
         (" src/utils/parse.ts  authorship:mixed  trust:0.71", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 PASS ", "JetBrains Mono", 11, EDDA, True),
         (" src/db/schema.ts  authorship:human  trust:0.95", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 BLOCK", "JetBrains Mono", 11, ERROR, True),
         (" src/ai/gen.ts  authorship:ai  trust:0.34  BOUNDARY", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 PASS ", "JetBrains Mono", 11, EDDA, True),
         (" src/core/engine.rs  authorship:human  trust:0.97", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 WARN ", "JetBrains Mono", 11, WARNING, True),
         (" src/lib/helpers.ts  authorship:ai  trust:0.68", "JetBrains Mono", 11, TEXT_MUTED, False)],
        [("\u2588 PASS ", "JetBrains Mono", 11, EDDA, True),
         (" src/config/policy.rego  authorship:human  trust:0.91", "JetBrains Mono", 11, TEXT_MUTED, False)],
    ]

    add_multiline_textbox(slide, right_pane_left + Inches(0.2), pane_top + Inches(0.5),
                          right_pane_w - Inches(0.4), Inches(3.5),
                          events, font_name="JetBrains Mono", font_size=11,
                          color=TEXT_PRIMARY, line_spacing=1.3)

    # Footer bar
    footer_top = frame_top + frame_h - Inches(0.6)
    add_rect(slide, frame_left, footer_top, frame_w, Inches(0.6),
             fill_color=SURFACE, border_color=STRUCTURE)
    add_textbox(slide, frame_left + Inches(0.3), footer_top + Inches(0.1),
                Inches(5), Inches(0.4),
                "[ \u25a0 ] e d d a c r a f t    v0.9.2-beta",
                font_name="JetBrains Mono", font_size=10, color=TEXT_MUTED)
    add_textbox(slide, frame_left + Inches(7), footer_top + Inches(0.1),
                Inches(5), Inches(0.4),
                "8 events  |  6 pass  |  1 warn  |  1 block  |  avg 47ms",
                font_name="JetBrains Mono", font_size=10, color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)

    add_footer(slide)
    add_presenter_notes(slide, "This is Anvil running in the terminal. The left pane shows the active policy -- the rules your team has defined. The right pane is the real-time signal interceptor -- it shows governance events as they happen. File saved, policy evaluated, architecture checked. The footer shows system logs. The product is built in Rust for performance and ships as a single binary. No runtime dependencies, no Docker containers, no cloud accounts required. It runs where your code runs. Every check is deterministic -- programmatic, mechanical, repeatable. No AI reviewing AI. The same input always produces the same output. This product plays in the exact space AI struggles with: precision. And it runs in under 50 milliseconds.")


def build_slide_07_market(prs):
    """Slide 7: Market — USD 21.5B market."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.4), CONTENT_W, Inches(0.7),
                "USD 21.5B market. USD 492M in AI governance alone.",
                font_name="JetBrains Mono", font_size=32, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Nested rectangles for TAM/SAM/SOM
    tam_left = Inches(1.5)
    tam_top = Inches(1.5)
    tam_w = Inches(10.3)
    tam_h = Inches(4.8)

    # TAM
    add_rect(slide, tam_left, tam_top, tam_w, tam_h,
             fill_color=SURFACE, border_color=STRUCTURE, border_width=Pt(2))
    add_textbox(slide, tam_left + Inches(0.3), tam_top + Inches(0.2),
                Inches(4), Inches(0.5), "USD 21.5B",
                font_name="JetBrains Mono", font_size=48, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, tam_left + Inches(0.3), tam_top + Inches(0.9),
                Inches(4), Inches(0.4), "TAM — AI code tools + AppSec + governance (2025)",
                font_name="Inter", font_size=14, color=TEXT_MUTED)

    # SAM
    sam_left = tam_left + Inches(1.5)
    sam_top = tam_top + Inches(1.5)
    sam_w = tam_w - Inches(3.0)
    sam_h = tam_h - Inches(2.2)

    add_rect(slide, sam_left, sam_top, sam_w, sam_h,
             fill_color=SURFACE, border_color=TEXT_MUTED, border_width=Pt(2))
    add_textbox(slide, sam_left + Inches(0.3), sam_top + Inches(0.2),
                Inches(4), Inches(0.5), "USD 1.5\u20132.0B",
                font_name="JetBrains Mono", font_size=36, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, sam_left + Inches(0.3), sam_top + Inches(0.7),
                Inches(5), Inches(0.4), "SAM — Governance + quality for AI-assisted development",
                font_name="Inter", font_size=13, color=TEXT_MUTED)

    # SOM
    som_left = sam_left + Inches(1.2)
    som_top = sam_top + Inches(1.2)
    som_w = sam_w - Inches(2.4)
    som_h = sam_h - Inches(1.5)

    add_rect(slide, som_left, som_top, som_w, som_h,
             fill_color=ANVIL_DIM, border_color=ANVIL, border_width=Pt(3))
    add_textbox(slide, som_left + Inches(0.3), som_top + Inches(0.2),
                Inches(4), Inches(0.5), "USD 50\u2013100M",
                font_name="JetBrains Mono", font_size=28, color=ANVIL, bold=True)
    add_textbox(slide, som_left + Inches(0.3), som_top + Inches(0.6),
                Inches(4), Inches(0.4), "SOM — Early adopter segment (Year 3)",
                font_name="Inter", font_size=13, color=TEXT_MUTED)

    # Side callouts
    add_textbox(slide, tam_left + Inches(6.5), tam_top + Inches(0.3),
                Inches(3.5), Inches(0.8),
                "3.4x more effective governance\nwith purpose-built platforms",
                font_name="Inter", font_size=14, color=ANVIL)

    add_textbox(slide, tam_left + Inches(6.5), tam_top + Inches(1.1),
                Inches(3.5), Inches(0.4),
                "AppSec: USD 13.6B, 22%+ CAGR",
                font_name="Inter", font_size=12, color=TEXT_MUTED)

    # Source
    add_textbox(slide, MARGIN, Inches(6.5), CONTENT_W, Inches(0.3),
                "Sources: Mordor Intelligence, Gartner (Feb 2026), derived analysis",
                font_name="Inter", font_size=9, color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)

    add_footer(slide)
    add_presenter_notes(slide, "The market sits at the intersection of three segments: AI code tools at 7.4 billion, application security testing at 13.6 billion, and AI governance platforms at 492 million and growing rapidly. The critical insight is that AI governance spend is not discretionary -- it is driven by regulatory deadlines and enterprise compliance requirements. Gartner's survey of 360 organisations found that those with governance platforms are 3.4 times more effective. This is not a nice-to-have; it is becoming infrastructure.")


def build_slide_08_competitive(prs):
    """Slide 8: Competitive Landscape — 2x2 matrix."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                "The only tool that is both deterministic and pre-commit",
                font_name="JetBrains Mono", font_size=30, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 2x2 matrix
    grid_left = Inches(2.5)
    grid_top = Inches(1.5)
    cell_w = Inches(4.0)
    cell_h = Inches(2.5)

    # Y-axis label
    add_textbox(slide, Inches(0.3), grid_top + Inches(0.5), Inches(1.8), Inches(0.5),
                "DETERMINISTIC",
                font_name="JetBrains Mono", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, Inches(0.3), grid_top + cell_h + Inches(0.5), Inches(1.8), Inches(0.5),
                "PROBABILISTIC",
                font_name="JetBrains Mono", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.RIGHT)

    # X-axis labels
    add_textbox(slide, grid_left, grid_top + 2 * cell_h + Inches(0.1), cell_w, Inches(0.4),
                "PRE-COMMIT",
                font_name="JetBrains Mono", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, grid_left + cell_w, grid_top + 2 * cell_h + Inches(0.1), cell_w, Inches(0.4),
                "POST-COMMIT",
                font_name="JetBrains Mono", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Grid lines
    add_line(slide, grid_left, grid_top, grid_left + 2 * cell_w, grid_top,
             color=STRUCTURE, width=Pt(2))
    add_line(slide, grid_left, grid_top + cell_h, grid_left + 2 * cell_w, grid_top + cell_h,
             color=STRUCTURE, width=Pt(2))
    add_line(slide, grid_left, grid_top + 2 * cell_h, grid_left + 2 * cell_w, grid_top + 2 * cell_h,
             color=STRUCTURE, width=Pt(2))
    add_line(slide, grid_left, grid_top, grid_left, grid_top + 2 * cell_h,
             color=STRUCTURE, width=Pt(2))
    add_line(slide, grid_left + cell_w, grid_top, grid_left + cell_w, grid_top + 2 * cell_h,
             color=STRUCTURE, width=Pt(2))
    add_line(slide, grid_left + 2 * cell_w, grid_top, grid_left + 2 * cell_w, grid_top + 2 * cell_h,
             color=STRUCTURE, width=Pt(2))

    # Top-left: ANVIL (highlighted)
    add_rect(slide, grid_left + Pt(2), grid_top + Pt(2),
             cell_w - Pt(4), cell_h - Pt(4),
             fill_color=ANVIL_DIM, border_color=ANVIL, border_width=Pt(3))
    add_textbox(slide, grid_left + Inches(0.3), grid_top + Inches(0.4),
                cell_w - Inches(0.6), Inches(0.5), "ANVIL",
                font_name="JetBrains Mono", font_size=28, color=ANVIL, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, grid_left + Inches(0.3), grid_top + Inches(1.0),
                cell_w - Inches(0.6), Inches(0.8),
                "Deterministic policy\nLine-level attribution\nPre-commit enforcement",
                font_name="Inter", font_size=13, color=TEXT_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    # Top-right: Static Analysis
    add_textbox(slide, grid_left + cell_w + Inches(0.3), grid_top + Inches(0.4),
                cell_w - Inches(0.6), Inches(0.5), "Static Analysis",
                font_name="Inter", font_size=20, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, grid_left + cell_w + Inches(0.3), grid_top + Inches(1.0),
                cell_w - Inches(0.6), Inches(0.8),
                "SonarQube, Semgrep, ESLint\nCodeQL, Snyk Code",
                font_name="Inter", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Bottom-left: Empty
    add_textbox(slide, grid_left + Inches(0.3), grid_top + cell_h + Inches(0.8),
                cell_w - Inches(0.6), Inches(0.5), "[empty quadrant]",
                font_name="Inter", font_size=14, color=STRUCTURE,
                alignment=PP_ALIGN.CENTER)

    # Bottom-right: AI Code Review
    add_textbox(slide, grid_left + cell_w + Inches(0.3), grid_top + cell_h + Inches(0.4),
                cell_w - Inches(0.6), Inches(0.5), "AI Code Review",
                font_name="Inter", font_size=20, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, grid_left + cell_w + Inches(0.3), grid_top + cell_h + Inches(1.0),
                cell_w - Inches(0.6), Inches(0.8),
                "CodeRabbit, Codacy AI\nGitHub Copilot Review\nAI Governance Platforms",
                font_name="Inter", font_size=13, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "This matrix shows the landscape on two axes: when governance happens -- pre-commit or post-commit -- and how it works -- deterministic policy or probabilistic AI. Static analysis tools like SonarQube and Semgrep are deterministic but post-commit. AI code review tools are probabilistic and post-commit. The top-left quadrant -- deterministic and pre-commit -- is empty. That is where Anvil sits. This is not a marginal improvement on an existing category. It is a structurally different approach. And moving into this quadrant requires fundamental re-architecture, not a feature addition.")


def build_slide_09_business_model(prs):
    """Slide 9: Business Model — Land with developers, expand with compliance."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                "Land with developers. Expand with compliance.",
                font_name="JetBrains Mono", font_size=34, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Horizontal expansion funnel (left side)
    funnel_left = Inches(0.5)
    funnel_top = Inches(1.4)

    stages = [
        ("CLI Install", "Free / community", TEXT_MUTED, Inches(1.2)),
        ("Team\nAdoption", "Per-seat subscription", ANVIL, Inches(1.5)),
        ("Enterprise\nPolicy", "Enterprise tier", ANVIL, Inches(1.8)),
        ("Compliance\nPacks", "Add-on revenue", EDDA, Inches(2.1)),
    ]

    stage_w = Inches(2.5)
    stage_gap = Inches(0.3)
    x = funnel_left

    for i, (label, revenue, color, h) in enumerate(stages):
        y_offset = (Inches(2.1) - h) // 2
        add_rect(slide, x, funnel_top + y_offset, stage_w, h,
                 fill_color=SURFACE, border_color=color, border_width=Pt(2))
        add_textbox(slide, x + Inches(0.2), funnel_top + y_offset + Inches(0.15),
                    stage_w - Inches(0.4), Inches(0.7), label,
                    font_name="JetBrains Mono", font_size=15, color=TEXT_PRIMARY,
                    bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), funnel_top + y_offset + h - Inches(0.5),
                    stage_w - Inches(0.4), Inches(0.4), revenue,
                    font_name="Inter", font_size=11, color=TEXT_MUTED,
                    alignment=PP_ALIGN.CENTER)

        # Arrow
        if i < 3:
            ax = x + stage_w
            ay = funnel_top + Inches(1.05)
            add_line(slide, ax + Inches(0.02), ay, ax + stage_gap - Inches(0.02), ay,
                     color=STRUCTURE, width=Pt(2))
        x += stage_w + stage_gap

    # Tier cards below
    tier_top = Inches(4.0)
    tier_h = Inches(2.5)
    tier_w = Inches(3.6)
    tier_gap = Inches(0.4)
    tier_start = Inches(1.3)

    tiers = [
        ("Community", "Open source core", STRUCTURE, [
            "CLI governance engine",
            "Core policy checks",
            "Local architecture graph",
        ]),
        ("Team", "Per-seat pricing", ANVIL, [
            "Everything in Community",
            "Policy packs (SOC 2, HIPAA)",
            "Team policy management",
            "Shared baselines",
        ]),
        ("Enterprise", "Custom pricing", EDDA, [
            "Everything in Team",
            "Centralised policy mgmt",
            "Audit dashboards + SSO",
            "Compliance reporting",
        ]),
    ]

    for i, (name, subtitle, color, features) in enumerate(tiers):
        tx = tier_start + i * (tier_w + tier_gap)
        add_rect(slide, tx, tier_top, tier_w, tier_h,
                 fill_color=SURFACE, border_color=color, border_width=Pt(2))

        # Top accent bar
        add_rect(slide, tx, tier_top, tier_w, Pt(4),
                 fill_color=color, border_color=color, border_width=Pt(0))

        add_textbox(slide, tx + Inches(0.2), tier_top + Inches(0.15),
                    tier_w - Inches(0.4), Inches(0.4), name,
                    font_name="JetBrains Mono", font_size=16, color=color, bold=True)
        add_textbox(slide, tx + Inches(0.2), tier_top + Inches(0.5),
                    tier_w - Inches(0.4), Inches(0.3), subtitle,
                    font_name="Inter", font_size=11, color=TEXT_MUTED)

        add_multiline_textbox(slide, tx + Inches(0.2), tier_top + Inches(0.85),
                              tier_w - Inches(0.4), Inches(1.5),
                              features, font_name="Inter", font_size=11,
                              color=TEXT_PRIMARY)

    add_footer(slide)
    add_presenter_notes(slide, "The go-to-market follows the developer tools playbook. A developer installs Anvil via CLI, configures a few policies, and sees immediate value -- governance events on every save. That is the land. The expand happens when compliance requirements arrive: the EU AI Act deadline, a SOC 2 audit, an enterprise customer asking about AI code governance. Policy packs -- pre-built rule sets for specific compliance frameworks -- become the expansion revenue. Enterprise features like centralised policy management and audit dashboards create the upsell.")


def build_slide_10_traction(prs):
    """Slide 10: Traction — Built what others are pitching."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                "Built what others are pitching",
                font_name="JetBrains Mono", font_size=36, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Three columns
    col_w = Inches(3.6)
    col_h = Inches(3.5)
    col_gap = Inches(0.35)
    col_start = Inches(0.85)
    col_top = Inches(1.2)

    columns = [
        ("BUILT TODAY", EDDA, "\u2713", [
            "Production Rust kernel (6 crates)",
            "50ms deterministic checks",
            "OPA/Rego policy engine",
            "Line-level authorship attribution",
            "Semantic graph + drift detection",
        ]),
        ("LAUNCH TRAJECTORY", ANVIL, "\u2192", [
            "5,000+ waitlist target",
            "10\u201315 pilot teams (5 engaged)",
            "Enterprise pipeline via Arkahna",
            "Developer influencer demos",
        ]),
        ("ECOSYSTEM", ANVIL, "\u2192", [
            "2 open source packages released",
            "Community building in progress",
            "IDE + CI/CD integrations planned",
        ]),
    ]

    for i, (header, color, prefix, items) in enumerate(columns):
        x = col_start + i * (col_w + col_gap)

        add_rect(slide, x, col_top, col_w, col_h,
                 fill_color=SURFACE, border_color=STRUCTURE, border_width=Pt(2))

        # Top accent
        add_rect(slide, x, col_top, col_w, Pt(3),
                 fill_color=color, border_color=color, border_width=Pt(0))

        # Header
        add_textbox(slide, x + Inches(0.2), col_top + Inches(0.2),
                    col_w - Inches(0.4), Inches(0.4), header,
                    font_name="JetBrains Mono", font_size=13, color=TEXT_MUTED, bold=True)

        # Items with prefixes
        lines = []
        for item in items:
            lines.append([
                (f"{prefix} ", "JetBrains Mono", 14, color, True),
                (item, "Inter", 14, TEXT_PRIMARY, False),
            ])

        add_multiline_textbox(slide, x + Inches(0.2), col_top + Inches(0.7),
                              col_w - Inches(0.4), col_h - Inches(1.0),
                              lines, font_name="Inter", font_size=14,
                              color=TEXT_PRIMARY, line_spacing=1.4)

    # Data callout below
    callout_y = col_top + col_h + Inches(0.3)
    add_rect(slide, col_start, callout_y, 3 * col_w + 2 * col_gap, Inches(1.2),
             fill_color=SURFACE, border_color=STRUCTURE)

    # Left accent
    add_rect(slide, col_start, callout_y, Pt(4), Inches(1.2),
             fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))

    add_textbox(slide, col_start + Inches(0.3), callout_y + Inches(0.15),
                3 * col_w + 2 * col_gap - Inches(0.6), Inches(0.5),
                "Not vibe-coded. Precision-engineered.",
                font_name="JetBrains Mono", font_size=24, color=ANVIL, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, col_start + Inches(0.3), callout_y + Inches(0.65),
                3 * col_w + 2 * col_gap - Inches(0.6), Inches(0.4),
                "Competitors in this category are raising on decks. Anvil is raising on a working product.",
                font_name="Inter", font_size=14, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "While other companies in the AI governance space are raising record rounds on pitch decks and prototypes, Anvil is a production-grade system. The Rust kernel, the policy engine, the semantic graph, the authorship attribution -- all built. Precision-engineered in a domain where AI struggles: deterministic analysis, sub-50-millisecond checks, repeatable results. We have 5 pilot teams today and developer influencers lined up to demo ahead of launch. We're targeting 5,000 on the waitlist and 10 to 15 pilot teams by the time we close this round. The product plays in the exact space AI fails at -- precision -- and that's the point.")


def build_slide_11_team(prs):
    """Slide 11: Team — 25 years building what enterprises buy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.3), CONTENT_W, Inches(0.6),
                "25 years building what enterprises buy",
                font_name="JetBrains Mono", font_size=34, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Left: portrait area (40%)
    portrait_left = MARGIN
    portrait_top = Inches(1.3)
    portrait_w = Inches(4.0)
    portrait_h = Inches(3.2)

    add_rect(slide, portrait_left, portrait_top, portrait_w, portrait_h,
             fill_color=SURFACE, border_color=STRUCTURE)

    # Monogram
    add_textbox(slide, portrait_left, portrait_top + Inches(0.8),
                portrait_w, Inches(1.2), "JB",
                font_name="JetBrains Mono", font_size=72, color=ANVIL,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, portrait_left, portrait_top + portrait_h - Inches(0.5),
                portrait_w, Inches(0.4), "Joshua Boys",
                font_name="JetBrains Mono", font_size=16, color=TEXT_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    # Right: credentials (60%)
    cred_left = portrait_left + portrait_w + Inches(0.5)
    cred_w = CONTENT_W - portrait_w - Inches(0.5)
    cred_top = Inches(1.3)

    add_textbox(slide, cred_left, cred_top, cred_w, Inches(0.5),
                "Joshua Boys",
                font_name="JetBrains Mono", font_size=28, color=TEXT_PRIMARY, bold=True)
    add_textbox(slide, cred_left, cred_top + Inches(0.5), cred_w, Inches(0.4),
                "Founder & CEO",
                font_name="Inter", font_size=20, color=TEXT_MUTED)

    credentials = [
        "Former Microsoft Azure Lead, Australia",
        "CEO of Arkahna \u2014 platform engineering for 100+ SaaS companies (5 years)",
        "25+ years building enterprise software, leading teams, shipping SaaS",
    ]

    cred_y = cred_top + Inches(1.1)
    for cred in credentials:
        add_textbox(slide, cred_left, cred_y, cred_w, Inches(0.4),
                    cred, font_name="Inter", font_size=16, color=TEXT_PRIMARY)
        cred_y += Inches(0.45)
        # Divider
        add_line(slide, cred_left, cred_y, cred_left + cred_w, cred_y,
                 color=STRUCTURE, width=Pt(1))
        cred_y += Inches(0.1)

    # Unfair advantage callout
    callout_top = cred_y + Inches(0.15)
    add_rect(slide, cred_left, callout_top, cred_w, Inches(0.7),
             fill_color=SURFACE, border_color=STRUCTURE)
    # Left accent
    add_rect(slide, cred_left, callout_top, Pt(5), Inches(0.7),
             fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))
    add_textbox(slide, cred_left + Inches(0.2), callout_top + Inches(0.1),
                cred_w - Inches(0.4), Inches(0.5),
                "Built governance tooling from inside the enterprise buying process \u2014 not from a research lab",
                font_name="Inter", font_size=14, color=TEXT_PRIMARY)

    # Bottom zone: advisory + capital efficiency
    bottom_top = Inches(5.0)

    # Advisory bench
    add_textbox(slide, MARGIN, bottom_top, Inches(6.5), Inches(0.7),
                "Advisory support across enterprise software, startup scaling, and large SaaS.\nFirst hires: engineering + enterprise-focused CRO \u2014 team scales with the raise.",
                font_name="Inter", font_size=14, color=TEXT_MUTED)

    # Capital efficiency metric
    metric_left = Inches(8.5)
    add_rect(slide, metric_left, bottom_top, Inches(4.0), Inches(0.8),
             fill_color=SURFACE, border_color=STRUCTURE)
    add_textbox(slide, metric_left + Inches(0.2), bottom_top + Inches(0.05),
                Inches(3.6), Inches(0.35), "\u00a30 raised \u2192 production product",
                font_name="JetBrains Mono", font_size=16, color=ANVIL, bold=True)
    add_textbox(slide, metric_left + Inches(0.2), bottom_top + Inches(0.4),
                Inches(3.6), Inches(0.35), "Capital efficient by design",
                font_name="Inter", font_size=12, color=TEXT_MUTED)

    # Closing beat
    add_textbox(slide, MARGIN, Inches(6.1), CONTENT_W, Inches(0.5),
                "Built to make you trust your AI more \u2014 not by asking you to trust ours.",
                font_name="JetBrains Mono", font_size=16, color=ANVIL,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "I've spent 25 years in enterprise software -- the last five as CEO of Arkahna, a platform engineering company that works with over 100 SaaS companies. I was the Azure lead in Australia for Microsoft. I know how enterprises buy developer tools, because I've been on both sides of that transaction. Anvil exists because I've watched AI coding tools arrive in my clients' organisations with zero governance. The advisory bench includes senior operators from enterprise, startups, and large SaaS. The team scales with this raise -- first hires are engineers and a CRO. Built to make you trust your AI more -- not by asking you to trust ours.")


def build_slide_12_the_ask(prs):
    """Slide 12: The Ask — Own the category before the window closes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.2), CONTENT_W, Inches(0.5),
                "Own the category before the window closes",
                font_name="JetBrains Mono", font_size=30, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    # Funding amount
    add_textbox(slide, MARGIN, Inches(0.7), CONTENT_W, Inches(1.0),
                "\u00a33\u20135M",
                font_name="JetBrains Mono", font_size=96, color=ANVIL, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(1.6), CONTENT_W, Inches(0.4),
                "\u00a315\u201325M pre-money valuation  |  Seed round",
                font_name="Inter", font_size=20, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Use of funds stacked bar
    bar_top = Inches(2.3)
    bar_h = Inches(0.7)
    bar_left = Inches(1.5)
    total_bar_w = Inches(10.3)

    segments = [
        ("Engineering", 0.40, ANVIL),
        ("Go-to-Market", 0.30, EDDA),
        ("Strategic Acq.", 0.20, TEXT_PRIMARY),
        ("Operations", 0.10, TEXT_MUTED),
    ]

    sx = bar_left
    for label, pct, color in segments:
        seg_w = int(total_bar_w * pct)
        add_rect(slide, sx, bar_top, seg_w, bar_h,
                 fill_color=color, border_color=color, border_width=Pt(0))

        # Label below
        add_textbox(slide, sx, bar_top + bar_h + Inches(0.05),
                    seg_w, Inches(0.5),
                    f"{label}\n{int(pct*100)}%",
                    font_name="JetBrains Mono", font_size=11, color=color,
                    alignment=PP_ALIGN.CENTER)
        sx += seg_w

    # Use of funds detail
    details_top = Inches(3.6)
    detail_items = [
        ("Engineering (~40%)", "3\u20134 hires. Rust kernel, platform layer, ecosystem integrations.", ANVIL),
        ("Go-to-Market (~30%)", "Enterprise CRO + developer advocacy. Arkahna\u2019s 100+ SaaS network.", EDDA),
        ("Strategic Acquisition (~20%)", "Platform engineering IP. Arm\u2019s-length, accelerating maturity.", TEXT_PRIMARY),
        ("Operations (~10%)", "Compliance certification, infrastructure.", TEXT_MUTED),
    ]

    for i, (title, desc, color) in enumerate(detail_items):
        y = details_top + i * Inches(0.5)
        add_textbox(slide, Inches(1.5), y, Inches(3.0), Inches(0.4),
                    title, font_name="JetBrains Mono", font_size=12, color=color, bold=True)
        add_textbox(slide, Inches(4.8), y, Inches(7.0), Inches(0.4),
                    desc, font_name="Inter", font_size=12, color=TEXT_MUTED)

    # Milestones timeline
    ml_top = Inches(5.8)
    ml_left = Inches(1.5)
    ml_right = Inches(11.8)
    ml_w = ml_right - ml_left

    add_line(slide, ml_left, ml_top + Inches(0.25), ml_right, ml_top + Inches(0.25),
             color=STRUCTURE, width=Pt(2))

    milestones = [
        ("Profitability", "Phase 1", EDDA, 0.0),
        ("50+ paying\nteams", "", ANVIL, 0.3),
        ("Aug 2026\nEU AI Act", "", ANVIL, 0.6),
        ("Phase 2\nKnowledge work", "", TEXT_MUTED, 0.95),
    ]

    for label, sublabel, color, pos in milestones:
        x = int(ml_left + ml_w * pos)
        dot_sz = Inches(0.15)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     x - dot_sz // 2, ml_top + Inches(0.25) - dot_sz // 2,
                                     dot_sz, dot_sz)
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

        add_textbox(slide, x - Inches(0.6), ml_top + Inches(0.45),
                    Inches(1.2), Inches(0.6), label,
                    font_name="Inter", font_size=11, color=color,
                    alignment=PP_ALIGN.CENTER, bold=(color == ANVIL))

    # Data callout
    add_textbox(slide, MARGIN, Inches(5.2), CONTENT_W, Inches(0.4),
                "EU AI Act enforcement: August 2026. 5 months to capture the compliance purchasing wave.",
                font_name="JetBrains Mono", font_size=14, color=ANVIL,
                alignment=PP_ALIGN.CENTER)

    # Phase 2 vision
    add_textbox(slide, MARGIN, Inches(6.8), CONTENT_W, Inches(0.3),
                "Phase 2: AI governance for all knowledge work \u2014 self-funded, no further dilution",
                font_name="Inter", font_size=14, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)
    add_presenter_notes(slide, "We're raising \u00a33\u20135M to own this category before the compliance window closes. 40% goes to engineering -- scaling the Rust kernel, building the platform layer, and ecosystem integrations. 30% to go-to-market -- an enterprise CRO and developer advocacy to drive bottom-up adoption. 20% to strategic acquisition of platform engineering IP -- proven infrastructure we can absorb rather than rebuild, accelerating our maturity by months. This round gets us to profitability on code governance. That's phase 1. Phase 2 is the bigger thesis: AI governance for all knowledge work -- legal, finance, operations -- starting from the beachhead where the pain is sharpest and the tooling is most mature. We reach phase 2 self-funded. No further dilution required.")


def build_slide_13_appendix(prs):
    """Slide 13: Appendix — Deep dive materials."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, MARGIN, Inches(0.6), CONTENT_W, Inches(0.7),
                "Deep dive materials",
                font_name="JetBrains Mono", font_size=40, color=TEXT_PRIMARY, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, MARGIN, Inches(1.3), CONTENT_W, Inches(0.4),
                "Reference material for Q&A and follow-up",
                font_name="Inter", font_size=18, color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER)

    # Section list
    sections = [
        ("[ = ]", "Technical Architecture", "Rust kernel, semantic graph, OPA integration, Ratatui TUI"),
        ("[ \u2261 ]", "Detailed Competitive Comparison", "Feature-by-feature across categories"),
        ("[ \u25a0 ]", "Regulatory Timeline", "EU AI Act, US frameworks, enterprise compliance deadlines"),
        ("[ = ]", "Financial Model Assumptions", "Pricing, adoption curve, revenue projections"),
        ("[ \u2261 ]", "Product Roadmap", "Near-term (6 months) and medium-term (18 months)"),
    ]

    card_left = Inches(2.0)
    card_w = Inches(9.3)
    y = Inches(2.2)

    for prefix, title, desc in sections:
        add_rect(slide, card_left, y, card_w, Inches(0.8),
                 fill_color=SURFACE, border_color=STRUCTURE)

        # Left accent
        add_rect(slide, card_left, y, Pt(4), Inches(0.8),
                 fill_color=ANVIL, border_color=ANVIL, border_width=Pt(0))

        add_textbox(slide, card_left + Inches(0.25), y + Inches(0.08),
                    Inches(0.7), Inches(0.3), prefix,
                    font_name="JetBrains Mono", font_size=14, color=ANVIL, bold=True)

        add_textbox(slide, card_left + Inches(0.9), y + Inches(0.05),
                    Inches(7.5), Inches(0.35), title,
                    font_name="JetBrains Mono", font_size=16, color=TEXT_PRIMARY, bold=True)

        add_textbox(slide, card_left + Inches(0.9), y + Inches(0.4),
                    Inches(7.5), Inches(0.35), desc,
                    font_name="Inter", font_size=13, color=TEXT_MUTED)

        y += Inches(0.95)

    add_footer(slide)
    add_presenter_notes(slide, "These slides are here for Q&A and follow-up. Each one provides the detailed evidence behind a claim in the main deck. The technical architecture slide shows how the Rust kernel, semantic graph, and OPA engine work together. The competitive comparison goes feature-by-feature. The regulatory timeline shows every deadline through 2030. Use these to answer deep questions with confidence.")


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()

    # Set 16:9 slide dimensions
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_slide_01_title(prs)
    build_slide_02_problem(prs)
    build_slide_03_why_now(prs)
    build_slide_04_solution(prs)
    build_slide_05_how_it_works(prs)
    build_slide_06_product(prs)
    build_slide_07_market(prs)
    build_slide_08_competitive(prs)
    build_slide_09_business_model(prs)
    build_slide_10_traction(prs)
    build_slide_11_team(prs)
    build_slide_12_the_ask(prs)
    build_slide_13_appendix(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                               "eddacraft-anvil-pitch-deck.pptx")
    prs.save(output_path)
    print(f"Deck saved to: {output_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
